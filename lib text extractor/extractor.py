import os
import pdfplumber
from docx import Document
from PIL import Image
import pytesseract
from pptx import Presentation
from openpyxl import load_workbook

# Define folders
input_folder = "input_files"
output_folder = "extracted_output"

# Create output folder if not exist
os.makedirs(output_folder, exist_ok=True)

# Extract from PDF
def extract_text_from_pdf(path):
    text = ""
    try:
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                text += page.extract_text() or ""
    except Exception as e:
        print(f"PDF Error [{path}]: {e}")
    return text.strip()

# Extract from DOCX
def extract_text_from_docx(path):
    try:
        doc = Document(path)
        return "\n".join(para.text for para in doc.paragraphs).strip()
    except Exception as e:
        print(f"DOCX Error [{path}]: {e}")
        return ""

# Extract from image
def extract_text_from_image(path):
    try:
        image = Image.open(path)
        return pytesseract.image_to_string(image).strip()
    except Exception as e:
        print(f"Image Error [{path}]: {e}")
        return ""

# Extract from PPTX
def extract_text_from_pptx(path):
    text = ""
    try:
        prs = Presentation(path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        print(f"PPTX Error [{path}]: {e}")
    return text.strip()

# Extract from XLSX
def extract_text_from_xlsx(path):
    text = ""
    try:
        wb = load_workbook(path, data_only=True)
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                text += row_text.strip() + "\n"
    except Exception as e:
        print(f"XLSX Error [{path}]: {e}")
    return text.strip()

# Loop through all files
for filename in os.listdir(input_folder):
    input_path = os.path.join(input_folder, filename)
    base_name, ext = os.path.splitext(filename)
    ext = ext.lower()
    output_path = os.path.join(output_folder, base_name + ".txt")

    extracted_text = ""

    if ext == ".pdf":
        extracted_text = extract_text_from_pdf(input_path)
    elif ext == ".docx":
        extracted_text = extract_text_from_docx(input_path)
    elif ext in [".jpg", ".jpeg", ".png"]:
        extracted_text = extract_text_from_image(input_path)
    elif ext == ".pptx":
        extracted_text = extract_text_from_pptx(input_path)
    elif ext == ".xlsx":
        extracted_text = extract_text_from_xlsx(input_path)
    else:
        print(f"❌ Skipped unsupported file: {filename}")
        continue

    if extracted_text:
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(extracted_text)
        print(f"✅ Extracted: {filename} → {output_path}")
    else:
        print(f"⚠️ No text found in: {filename}")